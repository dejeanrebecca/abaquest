import os
import re
import aspose.slides as slides
from pptx import Presentation

# Path to the docs folder relative to this script
DOCS_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "src", "docs")
OUTPUT_BASE_DIR = os.path.join(os.path.dirname(__file__), "output")

def get_quest_number(filename):
    match = re.search(r'AbaQuest-(\d+)', filename)
    if match:
        return match.group(1)
    return "Unknown"

import xml.etree.ElementTree as ET

def clean_svg_watermark(svg_path):
    tree = ET.parse(svg_path)
    root = tree.getroot()
    ns = {'svg': 'http://www.w3.org/2000/svg'}
    ET.register_namespace('', 'http://www.w3.org/2000/svg')
    ET.register_namespace('xlink', 'http://www.w3.org/1999/xlink')
    
    for g in root.findall('.//svg:g', ns):
        text_to_remove = []
        for text_elem in g.findall('./svg:text', ns):
            text_content = ''.join(text_elem.itertext())
            if 'Evaluation' in text_content or 'Aspose' in text_content:
                text_to_remove.append(text_elem)
                
        if text_to_remove:
            paths_to_remove = []
            for path_elem in g.findall('./svg:path', ns):
                if path_elem.get('stroke') == '#FF0000':
                    paths_to_remove.append(path_elem)
                    
            for t in text_to_remove:
                g.remove(t)
            for p in paths_to_remove:
                g.remove(p)
    
    # Write the cleaned data back
    tree.write(svg_path, xml_declaration=True, encoding='utf-8')

def convert_pptx_to_scenes(pptx_path, output_dir, quest_num):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    print(f"Processing Quest {quest_num}: {os.path.basename(pptx_path)}...")
    
    # 1. Export Slides as SVGs (The "Storyboard")
    print("  Exporting SVGs...")
    with slides.Presentation(pptx_path) as pres:
        for slide in pres.slides:
            # Name SVGs dynamically based on quest and scene
            svg_filename = f"quest{quest_num}_scene{slide.slide_number}.svg"
            svg_path = os.path.join(output_dir, svg_filename)
            with open(svg_path, "wb") as f:
                slide.write_as_svg(f)
            # Remove watermark from exported SVG
            clean_svg_watermark(svg_path)
                
    # 2. Extract Text and Notes (The "Script")
    print("  Extracting text and notes...")
    prs = Presentation(pptx_path)
    md_content = f"# Storyboard Specification: {os.path.basename(pptx_path)}\n\n"
    
    for i, slide in enumerate(prs.slides):
        scene_num = i + 1
        md_content += f"## Scene {scene_num}\n"
        md_content += f"![Scene {scene_num}](./quest{quest_num}_scene{scene_num}.svg)\n\n"
        
        # Extract Title and Body
        text_content = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        md_content += "### Content\n" + "\n".join([f"- {t}" for t in text_content if t.strip()]) + "\n\n"
        
        # Extract Speaker Notes (Your Narrator Script)
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if notes is not None and notes.strip():
            md_content += f"### Narrator Script (Cousin's Voice)\n> {notes.strip()}\n\n"
        md_content += "---\n\n"

    with open(os.path.join(output_dir, "specification.md"), "w", encoding="utf-8") as f:
        f.write(md_content)
    print(f"  Quest {quest_num} conversion complete! Check the '{output_dir}' folder.\n")

def main():
    if not os.path.exists(DOCS_DIR):
        print(f"Error: Could not find docs directory at {DOCS_DIR}")
        return
        
    pptx_files = [f for f in os.listdir(DOCS_DIR) if f.endswith('.pptx') and not f.startswith('~')]
    
    if not pptx_files:
        print("No .pptx files found in the docs directory.")
        return
        
    for pptx_file in pptx_files:
        pptx_path = os.path.join(DOCS_DIR, pptx_file)
        quest_num = get_quest_number(pptx_file)
        
        # Create a specific output folder for each quest
        output_dir = os.path.join(OUTPUT_BASE_DIR, f"Quest{quest_num}")
        
        convert_pptx_to_scenes(pptx_path, output_dir, quest_num)
        
    print("All conversions finished!")

if __name__ == "__main__":
    main()